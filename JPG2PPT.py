from pptx import Presentation
from pptx.util import Inches
import os
import tkinter as tk
from tkinter import filedialog, messagebox


def create_pptx_from_images(image_folder, output_pptx, pptx_width=Inches(10), pptx_height=Inches(7.5)):
    prs = Presentation()
    prs.slide_width = pptx_width
    prs.slide_height = pptx_height

    # Get all JPG images and sort them by name in descending order
    images = sorted([f for f in os.listdir(image_folder) if f.lower().endswith(".jpg")], reverse=True)

    for image in images:
        img_path = os.path.join(image_folder, image)
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide

        # Add the image to the slide, covering the full slide
        left = top = Inches(0)
        slide.shapes.add_picture(img_path, left, top, width=pptx_width, height=pptx_height)

    prs.save(output_pptx)
    messagebox.showinfo("Success", f"Presentation saved as {output_pptx}")


def browse_folder(entry_field):
    folder_selected = filedialog.askdirectory()
    entry_field.delete(0, tk.END)  # Clear existing content
    entry_field.insert(0, folder_selected)  # Insert selected folder


def convert_pptx(image_folder_entry, output_folder_entry):
    # Get the source folder path
    image_folder = image_folder_entry.get()
    if not image_folder:
        messagebox.showwarning("Warning", "Please select a source folder.")
        return

    # Get the destination folder path
    destination_folder = output_folder_entry.get()
    if not destination_folder:
        messagebox.showwarning("Warning", "Please select a destination folder.")
        return

    # Ask for a name for the pptx file
    output_pptx = os.path.join(destination_folder, "output_presentation.pptx")

    # Call the function to create PPTX
    create_pptx_from_images(image_folder, output_pptx, pptx_width=Inches(16), pptx_height=Inches(9))


# Setting up the GUI window
root = tk.Tk()
root.title("Image to PPTX Converter")

# Source folder input
source_label = tk.Label(root, text="Select Source Folder:")
source_label.pack(pady=5)

image_folder_entry = tk.Entry(root, width=50)
image_folder_entry.pack(pady=5)

source_button = tk.Button(root, text="Browse", command=lambda: browse_folder(image_folder_entry))
source_button.pack(pady=5)

# Destination folder input
destination_label = tk.Label(root, text="Select Destination Folder:")
destination_label.pack(pady=5)

output_folder_entry = tk.Entry(root, width=50)
output_folder_entry.pack(pady=5)

destination_button = tk.Button(root, text="Browse", command=lambda: browse_folder(output_folder_entry))
destination_button.pack(pady=5)

# Convert button
convert_button = tk.Button(root, text="Convert Images to PPTX",
                           command=lambda: convert_pptx(image_folder_entry, output_folder_entry), width=30, height=2)
convert_button.pack(pady=20)

root.mainloop()
